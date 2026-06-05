#utils

import json
from pathlib import Path

def load_system_paths():
    """Reads config.json and builds absolute file paths."""
    base_dir = Path(__file__).resolve().parent.parent
    config_path = base_dir / 'src' / 'config.json'
    
    with open(config_path, 'r') as file:
        config = json.load(file)
        
    input_path = base_dir / 'data' / config['input_folder'] / config['input_excel_name']
    template_path = base_dir / 'templates' / config['template_name']
    output_path = base_dir / 'output' / config['output_name']
    
    return input_path, template_path, output_path

def load_mapping():
    """Reads mapping.json and returns the list of dictionaries."""
    mapping_path = Path(__file__).resolve().parent / 'mapping.json'
    with open(mapping_path, 'r') as file:
        mapping_data = json.load(file)
    return mapping_data['mappings']



#core extractor
import openpyxl

def get_excel_value(excel_path, sheet_name, cell_ref):
    """Opens the Excel file, grabs a specific cell's value, and closes it."""
    try:
        # data_only=True ensures we get the calculated value, not the formula
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        sheet = wb[sheet_name]
        value = sheet[cell_ref].value
        wb.close()
        return value
    except Exception as e:
        print(f"Error extracting {cell_ref} from {sheet_name}: {e}")
        return None



#trans
def apply_transformations(ppt_placeholder, raw_value):
    """Applies specific text manipulations based on the target placeholder."""
    if raw_value is None:
        return ""
        
    raw_str = str(raw_value).strip()
    
    # ---------------------------------------------------------
    # Rule 1: Extract Month and Year from the Reporting Period
    # Example Input: "Reporting Period: 2026-MAY to 2026-MAY"
    # ---------------------------------------------------------
    if ppt_placeholder in ["{{MON}}", "{{YYYY}}"]:
        try:
            # 1. Split by 'to ' and take the last element -> "2026-MAY"
            last_part = raw_str.split("to ")[-1].strip()
            
            # 2. Split by '-' -> ["2026", "MAY"]
            year, month = last_part.split("-")
            
            # 3. Route to the correct placeholder
            if ppt_placeholder == "{{MON}}":
                return month
            elif ppt_placeholder == "{{YYYY}}":
                return year
                
        except Exception as e:
            print(f"Warning: Failed to parse date string '{raw_str}': {e}")
            return raw_str

    # ---------------------------------------------------------
    # Default Rule: Pass-through
    # If no special rules match, just return the raw string
    # ---------------------------------------------------------
    return raw_str



##ppt builder

from pptx import Presentation

class PPTBuilder:
    def __init__(self, template_path):
        """Initializes the builder by loading the PowerPoint template."""
        try:
            self.prs = Presentation(template_path)
            print(f"Loaded template: {template_path.name}")
        except Exception as e:
            print(f"Failed to load template: {e}")
            self.prs = None

    def update_placeholder(self, target_placeholder, new_value):
        """Scans all slides to find the placeholder and replaces it."""
        if not self.prs:
            return

        # Ensure we are passing a string to PowerPoint
        new_value_str = str(new_value)
        found = False

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Check if the placeholder is inside this shape's text
                if target_placeholder in shape.text:
                    # Replace the text
                    shape.text = shape.text.replace(target_placeholder, new_value_str)
                    found = True
                    print(f"Success: Updated '{target_placeholder}' -> '{new_value_str}'")

        if not found:
            print(f"Warning: Placeholder '{target_placeholder}' not found in the template.")

    def save_deck(self, output_path):
        """Saves the modified presentation to the output directory."""
        if self.prs:
            # Ensure the output directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(output_path)
            print(f"\n--- Deck successfully saved to: {output_path} ---")




#gen
# Import our custom modules
from utils_config import load_system_paths, load_mapping
from core_extractor import get_excel_value
from core_transformer import apply_transformations
from core_ppt_builder import PPTBuilder

def main():
    print("Starting SDLC Deck Automation...\n")

    # 1. Load system paths and mapping rules
    input_excel, template_ppt, output_ppt = load_system_paths()
    mappings = load_mapping()

    # 2. Initialize the PowerPoint Builder
    ppt = PPTBuilder(template_ppt)
    if not ppt.prs:
        return # Stop if the template failed to load

    # 3. Process each mapping rule one by one
    print("\nProcessing Data Mappings:")
    for rule in mappings:
        placeholder = rule.get("ppt_placeholder")
        sheet = rule.get("excel_sheet")
        cell = rule.get("excel_cell")

        # Step A: Extract
        raw_value = get_excel_value(input_excel, sheet, cell)
        
        # Step B: Transform
        clean_value = apply_transformations(placeholder, raw_value)
        
        # Step C: Load (Update PPT)
        ppt.update_placeholder(placeholder, clean_value)

    # 4. Save the final output
    ppt.save_deck(output_ppt)

if __name__ == "__main__":
    main()
